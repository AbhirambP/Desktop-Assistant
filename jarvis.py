import sys
import speech_recognition as sr
import pyttsx3
import pywhatkit
import datetime
import wikipedia
import pyjokes
import webbrowser
import time
import os
import tkinter as tk
import cv2
from requests import get
import psutil
import pyautogui
import PyPDF2
from bs4 import BeautifulSoup
from PyQt5.QtGui import *
from PyQt5.QtCore import *
from PyQt5.QtWidgets import *
from pywikihow import search_wikihow
from pytube import YouTube
import qrcode
import string
import random
import speedtest_cli
from GoogleNews import GoogleNews
import pandas as pd
import winshell
import shutil
import comtypes.client
from pptx import Presentation
import pdf2docx
from docx2pdf import convert as docx2pdf_convert
import pywhatkit as kit

#Set our engine to "Pyttsx3" which is used for text to speech in Python 
#and sapi5 is Microsoft speech application platform interface 

engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice',voices[0].id) 


class Jarvis():
    def __init__(self):
        self.engine = pyttsx3.init()
        self.command = ""

    def take_command(self):
        try:
            listener = sr.Recognizer()
            with sr.Microphone() as source:
                print('Listening....')
                listener.pause_threshold = 1
                voice = listener.listen(source, timeout=4, phrase_time_limit=7)
                print("Recognizing...")
                command1 = listener.recognize_google(voice, language='en-in')
                command1 = command1.lower()
                if 'jarvis' in command1:
                    command1 = command1.replace('jarvis', '')
                return command1
        except sr.WaitTimeoutError:
            print("Listening timed out while waiting for phrase to start.")
            return 'None'
        except sr.UnknownValueError:
            print("Google Speech Recognition could not understand the audio.")
            return 'None'
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return 'None'
        except Exception as e:
            print(f"An error occurred: {e}")
            return 'None'
        
    def speak(self,audio):
       engine.say(audio)
       engine.runAndWait()

    def take_command_with_fallback(self):
        command = self.take_command()
        if command == 'None':
            print("No voice input detected. Please type your query:")
            self.speak("No voice input detected. Please type your query")
            command = input("Type your query: ").lower()
        return command
    
    def run_jarvis(self,output_folder):
        self.wishme()
        self.speak('Hello sir I am jarvis your assistant. please tell me how can i help you')
        while True:
            command = self.take_command_with_fallback() #Every time taking command after a task is done
            print(f"Command: {command}")
            
            if 'play a song' in command or 'youtube' in command or "download a song" in command or "download song" in command :
                #commands for opening youtube, playing a song in youtube, and download a song in youtube
                self.yt(command) #function is from line 555
            #Interaction commands with JARVIS
            elif 'your age' in command or 'are you single'in command or 'are you there' in command or 'tell me something' in command or 'thank you' in command or 'in your free time' in command or 'i love you' in command or 'can you hear me' in command or 'do you ever get tired' in command:
                self.Fun(command)
            elif 'time' in command : 
                self.Clock_time(command)
            elif 'hi' in command and len(command)==2 or 'hai' in command or 'hey' in command and len(command)==3 or 'hello' in command and len(command)==5 :
                self.comum(command)
            elif 'what can you do' in command or 'your name' in command or 'my name' in command or 'college name' in command:
                self.Fun(command)
            elif 'joke'in command or 'date' in command:
                self.Fun(command)
            #schedule commands for remembering you what is the planns of the day
            elif "college time table" in command or "schedule" in command :
                self.shedule() #function is present from 407
            #It will tell the day Eg : Today is wednesday
            elif "today" in command :
                day = self.Cal_day()
                self.speak("Today is "+day)
            #commad for opening any weekly meeting links
            #Eg: I have kept a meeting my amFOSS club 
            #Note: the given link is fake!!
            elif "meeting" in command:
                self.speak("Ok sir opening meeet")
                webbrowser.open("https://meeting/")
            #command if you don't want the JARVIS to speak until for a certain time
            #Note: I can be silent for max of 10mins
            # Eg: JARVIS keep quiet for 5 minutes 
            elif 'silence' in command or 'silent' in command or 'keep quiet' in command or 'wait for' in command :
                self.silenceTime(command)
            #Command for opening your social media accounts in webrowser
            #Eg : JARVIS open facebook (or) JARVIS open social media facebook 
            elif 'facebook' in command or 'whatsapp' in command or 'instagram' in command or 'twitter' in command or 'discord' in command or 'social media' in command :
                self.social(command)
            #command for opening your OTT platform accounts
            #Eg: open hotstart
            elif 'hotstar' in command or 'prime' in command or 'netflix' in command :
                self.OTT(command)
            #Command for opening your online classes links
            elif 'online classes'in command :
                self.OnlineClasses(command)
            #command for opeing college websites
            elif 'open teams'in command or 'mlrit portal'in command or 'open sharepoint'in command or'open outlook'in command or'lms'in command :
                self.college(command)
            #command to search for something in wikipedia
            #Eg: what is meant by python in wikipedia (or) search for "_something_" in wikipedia
            elif 'wikipedia' in command or 'what is meant by' in command or 'tell me about' in command or 'who the heck is' in command :
                self.B_S(command)
            #command for opening your browsers and search for information in google
            elif 'open google'in command or 'open edge'in command :
                self.brows(command)
            #command to open your google applications
            elif 'open gmail'in command or'open maps'in command or'open calender'in command or'open documents'in command or'open spredsheet'in command or'open images'in command or'open drive'in command or'open news' in command :
                self.Google_Apps(command)
            #command to open your open-source accounts
            #you can add other if you have
            elif 'open github'in command or 'open gitlab'in command :
                self.open_source(command)
            #commands to open presentaion makeing tools like CANVA and GOOGLE SLIDES
            elif 'slides' in command or 'canva'in command :
                self.edit(command)
            #Command to open desktop applications
            #It can open : calculator, notepad,paint, teams(aka online classes), discord, spotify, ltspice,vscode(aka editor), steam, VLC media player
            elif 'open calculator'in command or 'open notepad'in command or 'open paint'in command or 'open online classes'in command or 'open discord'in command or 'open editor'in command or 'open spotify'in command or 'open media player'in command :
                self.OpenApp(command)
            #Command to close desktop applications
            #It can close : caliculator, notepad,paint, discord, spotify, ltspice,vscode(aka editor), steam, VLC media player
            elif 'close calculator'in command or 'close notepad'in command or 'close paint'in command or 'close discord'in command or 'close editor'in command or 'close spotify'in command or 'close media player'in command:
                self.CloseApp(command)
            #command for opening shopping websites 
            #NOTE: you can add as many websites
            elif 'flipkart'in command or 'amazon'in command :
                self.shopping(command)
            #command for asking your current location
            elif 'where i am' in command or 'where we are' in command:
                self.locaiton()
            #command for opening command prompt 
            #Eg: jarvis open command prompt
            elif 'command prompt'in command :
                self.speak('Opening command prompt')
                os.system('start cmd')
            elif 'take a note' in command or 'write in notepad' in command:
                self.speak("opening notepad")
                self.open_notepad_and_dictate(command)
            #Command for opening taking screenshot
            #Eg: jarvis take a screenshot
            elif 'take screenshot' in command or 'screenshot' in command or"take a screenshot" in command :
                self.scshot()
            #command for searching for a procedure how to do something
            #Eg:jarvis activate mod
            #   jarvis How to make a cake (or) jarvis how to convert int to string in programming 
            elif "activate mod" in command:
                self.How()
            #command for increaing the volume in the system
            #Eg: jarvis increase volume
            elif "volume up" in command or "increase volume" in self.command :
                pyautogui.press("volumeup")
                self.talk('volume increased')
            #command for decreaseing the volume in the system
            #Eg: jarvis decrease volume
            elif "volume down" in command or "decrease volume" in command:
                pyautogui.press("volumedown")
                self.speak('volume decreased')
            #Command to mute the system sound
            #Eg: jarvis mute the sound
            elif "volume mute" in command or "mute the sound" in command or "turn off the volume" in command :
                pyautogui.press("volumemute")
                self.speak('volume muted')
            #command for opening your mobile camera the description for using this is in the README file
            #Eg: Jarvis open mobile camera
            elif "open mobile cam" in command :
                self.Mobilecamra()
            #command for opening your webcamera
            #Eg: jarvis open webcamera
            elif 'open cam'in command or "open camera" in command:
                self.open_camera()

            elif "display all the contacts" in command:
                self.Display()
            #Command for checking covid status in India
            #Eg: jarvis check covid (or) corona status
            elif "covid" in command or  "corona" in command :
                self.speak("Boss which state covid 19 status do you want to check")
                s = self.take_command_with_fallback()
                self.Covid(s)

            #command for playing a dowloaded mp3 song in which is present in your system
            #Eg: Jarvis play music
            elif 'music' in command:
                try:
                    music_dir = 'E:\\music' #change the song path directory if you have songs in other directory
                    songs = os.listdir(music_dir)
                    for song in songs:
                        if song.endswith('.mp3'):
                            os.startfile(os.path.join(music_dir, song))
                except:
                    self.speak("Boss an unexpected error occured")
            #command for knowing your system IP address
            #Eg: jarvis check my ip address
            elif 'ip address' in command:
                ip = get('https://api.ipify.org').text
                print(f"your IP address is {ip}")
                self.speak(f"your IP address is {ip}")

            #command for checking the temperature in surroundings
            #jarvis check the surroundings temperature
            elif "temperature" in command:
                self.temperature()
            #Command to generate the qr codes
            elif "create a qr code" in command:
                self.qrCodeGenerator()
            #command for checking internet speed
            #Eg: jarvis check my internet speed
            elif "internet speed" in command:
                self.InternetSpeed()
            #command to make the jarvis sleep
            #Eg: jarvis you can sleep now
            elif "you can sleep" in command or "sleep now" in command :
                self.speak  ("Okay boss, I am going to sleep you can call me anytime.")
                break
            #command for waking the jarvis from sleep
            #jarvis wake up
            elif "wake up" in command or "get up" in command :
                self.speak("yes sir , how can i help you")
            #command for exiting jarvis from the program
            #Eg: jarvis goodbye
            elif "goodbye" in command or "get lost" in command :
                self.speak("Thanks for using me boss, have a good day")
                sys.exit()
            #command for knowing about your system condition
            #Eg: jarvis what is the system condition
            elif 'system condition' in command or 'condition of the system' in command:
                self.speak("checking the system condition")
                self.condition()
            #command for knowing the latest news
            #Eg: jarvis tell me the news
            elif 'tell me news' in command or "the news" in command or "todays news" in command :
                self.speak("Please wait boss, featching the latest news")
                self.news()
            #command for shutting down the system
            #Eg: jarvis shutdown the system
            elif 'shutdown the system' in command or 'down the system' in command:
                self.speak("shutting down the system in 10 seconds")
                time.sleep(10)
                os.system("shutdown /s /t 5")
            #command for restarting the system
            #Eg: jarvis restart the system
            elif 'restart the system' in command:
                self.speak("restarting the system in 10 seconds")
                time.sleep(10)
                os.system("shutdown /r /t 5")
            #command for make the system sleep
            #Eg: jarvis sleep the system
            elif 'sleep the system' in command:
                self.speak("system is going to sleep")
                os.system("rundll32.exe powrprof.dll, SetSuspendState 0,1,0")
            #command for performing windows actions
            #Eg: like switching tabs , jarvis open newtab , new window
            elif "switch window" in command or "open new tab" in command or "open new window" in command or "go to next window" in command:
                self.windows_manager(command)

            elif "keep windows side by side" in command or "go to previous window" in command or 'open history' in command or 'open downloads' in command:
                self.windows_manager(command)
            #command for closing all tabs and clearing browsing history
            elif 'clear browsing history' in command or 'close all tabs' in command:
                self.windows_manager(command)
            #command for performing folder operations
            #Eg: to open folder , delete folder, creating , renameing folders
            elif 'open folder' in command or 'delete folder' in command or 'create folder' in command or 'rename folder' in command or 'close folder' in command:
                self.folder_management(command)
            

            elif "open pdf" in command or "convert pdf to word" in command or "convert pdf to ppt" in command:
                self.documents_manager(command , output_folder , self.search_file)
            elif "convert word to pdf" in command or "convert ppt to pdf" in command :
                self.documents_manager(command, output_folder , self.search_file)

            elif "send message" in command :
                self.whatsapp(command)

    def Intro(self):
        while True:                
            self.permission = self.take_command_with_fallback()
            print(self.permission)
            if ("wake up" in self.permission) or ("get up" in self.permission):
                self.run_jarvis()
            elif ("goodbye" in self.permission) or ("get lost" in self.permission):
                self.talk("Thanks for using me boss, have a good day")
                sys.exit()

    
    def wishme(self):

        hour = datetime.datetime.now().hour
        if hour >= 0 and hour < 12:
            self.speak("Good Morning!")
        elif hour >= 12 and hour < 18:
            self.speak("Good Afternoon!")
        elif hour >= 18 and hour < 24:
            self.speak("Good Evening!")
        else:
            self.speak("Good Night Sir, See You Tomorrow")

    #Weather forecast
    def temperature(self):
        IP_Address = get('https://api.ipify.org').text
        url = 'https://get.geojs.io/v1/ip/geo/'+IP_Address+'.json'
        geo_reqeust = get(url)
        geo_data = geo_reqeust.json()
        city = geo_data['city']
        search = f"temperature in {city}"
        url_1 = f"https://www.google.com/search?q={search}"
        r = get(url_1)
        data = BeautifulSoup(r.text,"html.parser")
        temp = data.find("div",class_="BNeawe").text
        self.speak(f"current {search} is {temp}")

    #qrCodeGenerator
    def qrCodeGenerator(self):
        self.speak(f"Boss enter the text/link that you want to keep in the qr code")
        input_Text_link = input("Enter the Text/Link : ")
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=15,
            border=4,
        )
        QRfile_name = (str(datetime.datetime.now())).replace(" ","-")
        QRfile_name = QRfile_name.replace(":","-")
        QRfile_name = QRfile_name.replace(".","-")
        QRfile_name = QRfile_name+"-QR.png"
        qr.add_data(input_Text_link)
        qr.make(fit=True)

        img = qr.make_image(fill_color="black", back_color="white")
        img.save(f"QRCodes\{QRfile_name}")
        self.speak(f"Boss the qr code has been generated")
    
    #password generator
    def password():
        char = string.ascii_letters + string.digits
        ret = ''.join(random.choice(char) for x in range(0, 10))
        print(ret)

    #notpad writing

    def open_notepad_and_dictate(self,commmand):
        npath = "C:\\Windows\\system32\\notepad.exe"
        self.speak("Opening notepad")
        print("Opening notepad...")
        os.startfile(npath)
        time.sleep(2)  # Give Notepad time to open
    
        self.speak("Start dictating your notes.")
        while True:
            command = self.take_command_with_fallback().lower()
            if 'stop dictation' in command or 'stop' in command:
                self.speak("Stopping dictation.")
                break
            pyautogui.typewrite(command + '\n', interval=0.05)

    def create_folder(self,folder_name):
        if not folder_name:
            self.speak("Folder name is empty. Please specify a valid folder name.")
            return
        try:
            os.makedirs(folder_name, exist_ok=True)
            self.speak(f"Folder '{folder_name}' created successfully.")
        except Exception as e:
            self.speak(f"An error occurred while creating the folder: {str(e)}")

    def find_folder(folder_name, search_path):
        for root, dirs, files in os.walk(search_path):
            if folder_name in dirs:
                return os.path.join(root, folder_name)
        return None


    def open_folder(self,folder_name):
        if not folder_name:
            self.speak("Folder name is empty. Please specify a valid folder name.")
            return
    
        search_paths = [os.path.expanduser("~"), "C:\\"] if os.name == 'nt' else [os.path.expanduser("~")]
        folder_path = None
    
        for path in search_paths:
            folder_path = self.find_folder(folder_name, path)
            if folder_path:
                break

        if folder_path:
            try:
                if os.name == 'posix':  # macOS or Linux
                    subprocess.call(['open', folder_path])
                elif os.name == 'nt':  # Windows
                    os.startfile(folder_path)
                self.speak(f"Folder '{folder_name}' opened successfully.")
            except Exception as e:
                self.speak(f"An error occurred while opening the folder: {str(e)}")
        else:
           self.speak(f"Folder '{folder_name}' not found on the system.")

    def rename_folder(self,old_name, new_name):
        if not old_name or not new_name:
            self.speak("Folder names cannot be empty. Please specify valid folder names.")
            return
        try:
            if os.path.isdir(old_name):
                os.rename(old_name, new_name)
                self.speak(f"Folder '{old_name}' renamed to '{new_name}' successfully.")
            else:
                self.speak(f"Folder '{old_name}' does not exist.")
        except Exception as e:
            self.speak(f"An error occurred while renaming the folder: {str(e)}")

    def delete_folder_on_desktop(self,folder_name):
        desktop_path = winshell.desktop()
        folder_path = os.path.join(desktop_path, folder_name)
        if os.path.isdir(folder_path):
            shutil.rmtree(folder_path)
            self.speak(f"Folder '{folder_name}' deleted from desktop.")
        else:
            self.speak(f"Folder '{folder_name}' does not exist on desktop.")
            
    #whatsapp
    def whatsapp(self,command):
        try:
            command = command.replace('send a message to','')
            command = command.strip()
            name,numberID,F = self.SearchCont(command)
            if F:
                print(numberID)
                self.speak(f'Boss, what message do you want to send to {name}')
                message = self.take_command_with_fallback()
                hour = int(datetime.datetime.now().hour)
                min = int(datetime.datetime.now().minute)
                print(hour,min)
                if "group" in command:
                    kit.sendwhatmsg_to_group(numberID,message,int(hour),int(min)+1)
                else:
                    kit.sendwhatmsg(numberID,message,int(hour),int(min)+1)
                self.speak("Boss message have been sent")
            if F==False:
                self.speak(f'Boss, the name not found in our data base, shall I add the contact')
                AddOrNot = self.take_command_with_fallback()
                print(AddOrNot)
                if ("yes" in AddOrNot) or ("add" in AddOrNot) or ("yeah" in AddOrNot) or ("yah" in AddOrNot):
                    self.AddContact()
                elif("no" in AddOrNot):
                    self.speak('Ok Boss')
        except:
            print("Error occured, please try again")


    def Mobilecamra(self):
        import urllib.request
        import numpy as np
        try:
            self.speak(f" openinging mobile camera")
            URL = "http://_IP_Webcam_IP_address_/shot.jpg" #Discription for this is available in the README file
            while True:
                imag_arr = np.array(bytearray(urllib.request.urlopen(URL).read()),dtype=np.uint8)
                img = cv2.imdecode(imag_arr,-1)
                cv2.imshow('IPWebcam',img)
                q = cv2.waitKey(1)
                if q == ord("q"):
                    self.speak(f" closing mobile camera")
                    break
            cv2.destroyAllWindows()
        except Exception as e:
            print("Some error occured")

    #Web camera
    #NOTE to exit from the web camera press "ESC" key 
    def open_camera():
    # Open the default camera (usually the first one, index 0)
        cap = cv2.VideoCapture(0)

        if not cap.isOpened():
            print("Error: Could not open camera.")
            return

        while True:
        # Capture frame-by-frame
            ret, frame = cap.read()

        # Display the resulting frame
            cv2.imshow('Camera Feed', frame)

        # Exit loop if 'Esc' key is pressed
            if cv2.waitKey(1) == 27:
                break

    # Release the camera and close all OpenCV windows
        cap.release()
        cv2.destroyAllWindows()


    
    #Add contacts
    def AddContact(self):
        self.speak(f'Boss, Enter the contact details')
        name = input("Enter the name :").lower()
        number = input("Enter the number :")
        NumberFormat = f'"{name}":"+91{number}"'
        ContFile = open("Contacts.txt", "a") 
        ContFile.write(f"{NumberFormat}\n")
        ContFile.close()
        self.speak(f'Boss, Contact Saved Successfully')

    #Search Contact
    def SearchCont(self,name):
        with open("Contacts.txt","r") as ContactsFile:
            for line in ContactsFile:
                if name in line:
                    print("Name Match Found")
                    s = line.split("\"")
                    return s[1],s[3],True
        return 0,0,False
    
    #Display all contacts
    def Display(self):
        ContactsFile = open("Contacts.txt","r")
        count=0
        for line in ContactsFile:
            count+=1
        ContactsFile.close()
        ContactsFile = open("Contacts.txt","r")
        self.speak(f"Boss displaying the {count} contacts stored in our data base")    
        for line in ContactsFile:
            s = line.split("\"")
            print("Name: "+s[1])
            print("Number: "+s[3])
        ContactsFile.close()

    #search contact
    def NameIntheContDataBase(self,command):
        line = command
        line = line.split("number in contacts")[0]
        if("tell me" in line):
            name = line.split("tell me")[1]
            name = name.strip()
        else:
            name= line.strip()
        name,number,bo = self.SearchCont(name)
        if bo:
            print(f"Contact Match Found in our data base with {name} and the mboile number is {number}")
            self.speak(f"Boss Contact Match Found in our data base with {name} and the mboile number is {number}")
        else:
            self.speak("Boss the name not found in our data base, shall I add the contact")
            AddOrNot = self.take_command_with_fallback()
            print(AddOrNot)
            if ("yes add it" in AddOrNot)or ("yeah" in AddOrNot) or ("yah" in AddOrNot):
                self.AddContact()
                self.speak(f'Boss, Contact Saved Successfully')
            elif("no" in AddOrNot) or ("don't add" in AddOrNot):
                self.speak('Ok Boss')

    #Internet spped
    def InternetSpeed(self):
        self.speak("Wait a few seconds boss, checking your internet speed")
        st = speedtest_cli.Speedtest()
        dl = st.download()
        dl = dl/(1000000) #converting bytes to megabytes
        up = st.upload()
        up = up/(1000000)
        print(dl,up)
        self.speak(f"Boss, we have {dl} megabytes per second downloading speed and {up} megabytes per second uploading speed")
        
    #Search for a process how to do
    def How(self):
        self.speak("How to do mode is is activated")
        while True:
            self.speak("Please tell me what you want to know")
            how = self.take_command_with_fallback()
            try:
                if ("exit" in how) or("close" in how):
                    self.speak("Ok sir how to mode is closed")
                    break
                else:
                    max_result=1
                    how_to = search_wikihow(how,max_result)
                    assert len(how_to) == 1
                    how_to[0].print()
                    self.speak(how_to[0].summary)
            except Exception as e:
                self.speak("Sorry sir, I am not able to find this")

    #Communication commands
    def comum(self,command):
        print(command)
        if ('hi'in command) or('hai'in command) or ('hey'in command) or ('hello' in command) :
            self.speak("Hello boss what can I help for u")
        else :
            self.No_result_found()

    def folder_management(self,command):
        print(command)
        if 'create folder' in command:
            folder_name = command.replace('create folder', '').strip()
            if folder_name:
                self.create_folder(folder_name)
            else:
                self.speak("Please specify a folder name.")

        elif 'open folder' in command:
            folder_name = command.replace('open folder', '').strip()
            if folder_name:
                self.open_folder(folder_name)
            else:
                self.speak("Please specify a folder name.")

        elif 'rename folder' in command:
            self.speak("Please say the current folder name.")
            old_name = self.take_command_with_fallback()
            if old_name:
                self.speak("Please say the new folder name.")
                new_name = self.take_command_with_fallback()
                if new_name:
                    self.rename_folder(old_name.strip(), new_name.strip())

        elif 'delete folder' in command:
            folder_name = command.split("delete folder ")[-1]

        elif "close folder" in command:
    # Simulate Alt + F4 keyboard shortcut to close current window
            pyautogui.hotkey('alt', 'f4')
            print("Folder window closed.")
            self.speak("Folder window closed")


    #Fun commands to interact with jarvis
    def Fun(self,command):
        print(command)
        if 'your name' in command:
            self.speak("My name is jarvis")
        elif 'my name' in command:
            self.speak("your name is abhiram")
        elif 'college name' in command:
            self.speak("you are studing in MLR institution of technology, with batchelor in Computer Science and Artificail Intelligence and Machine Learning") 
        elif 'what can you do' in command:
            self.speak("I talk with you until you want to stop, I can say time, open your social media accounts,your open source accounts, open google browser,and I can also open your college websites, I can search for some thing in google and I can tell jokes")
        elif 'your age' in command:
            self.speak("I am very young that u")
        elif 'can we have a date' in command:
            self.speak('Sorry not intreseted, I am having headache, we will catch up some other time')
        elif 'are you single' in command:
            self.speak('No, I am in a relationship with wifi')
        elif 'joke' in command:
            self.speak(pyjokes.get_joke())
        elif 'are you there' in command:
            self.speak('Yes sir I am here')
        elif 'tell me something' in command:
            self.speak('I don\'t have much to say, you only tell me someting i will give you the company')
        elif 'thank you' in command:
            self.speak('I am here to help you..., your welcome')
        elif 'in your free time' in self.command:
            self.speak('I will be listening to all your words')
        elif 'i love you' in command:
            self.speak('I love you too boss')
        elif 'can you hear me' in command:
            self.speak('Yes sir, I can hear you')
        else :
            self.No_result_found()

    #Social media accounts commands
    def social(self,command):
        print(command)
        if 'facebook' in command:
            self.speak('opening your facebook')
            webbrowser.open('https://www.facebook.com/')
        elif 'whatsapp' in command:
            self.speak('opening your whatsapp')
            webbrowser.open('https://web.whatsapp.com/')
        elif 'instagram' in command:
            self.speak('opening your instagram')
            webbrowser.open('https://www.instagram.com/')
        elif 'twitter' in command:
            self.speak('opening your twitter')
            webbrowser.open('https://twitter.com/Suj8_116')
        elif 'discord' in command:
            self.speak('opening your discord')
            webbrowser.open('https://discord.com/channels/@me')
        else :
            self.No_result_found()
        
    #clock commands
    def Clock_time(self,command):
        print(command)
        time = datetime.datetime.now().strftime('%I:%M %p')
        print(time)
        self.speak("Current time is "+time)
    
    #calender day
    def Cal_day(self):
        day = datetime.datetime.today().weekday() + 1
        Day_dict = {1: 'Monday', 2: 'Tuesday', 3: 'Wednesday',4: 'Thursday', 5: 'Friday', 6: 'Saturday',7: 'Sunday'}
        if day in Day_dict.keys():
            day_of_the_week = Day_dict[day]
            print(day_of_the_week)
        
        return day_of_the_week

    #shedule function for remembering todays plans
    #NOTE For example I have declared my college timetable you can declare anything you want
    def shedule(self):
        day = self.Cal_day().lower()
        self.speak("Boss today's shedule is")
        Week = {"monday" : "Boss from 9:00 to 9:50 you have Cultural class, from 10:00 to 11:50 you have mechanics class, from 12:00 to 2:00 you have brake, and today you have sensors lab from 2:00",
        "tuesday" : "Boss from 9:00 to 9:50 you have English class, from 10:00 to 10:50 you have break,from 11:00 to 12:50 you have ELectrical class, from 1:00 to 2:00 you have brake, and today you have biology lab from 2:00",
        "wednesday" : "Boss today you have a full day of classes from 9:00 to 10:50 you have Data structures class, from 11:00 to 11:50 you have mechanics class, from 12:00 to 12:50 you have cultural class, from 1:00 to 2:00 you have brake, and today you have Data structures lab from 2:00",
        "thrusday" : "Boss today you have a full day of classes from 9:00 to 10:50 you have Maths class, from 11:00 to 12:50 you have sensors class, from 1:00 to 2:00 you have brake, and today you have english lab from 2:00",
        "friday" : "Boss today you have a full day of classes from 9:00 to 9:50 you have Biology class, from 10:00 to 10:50 you have data structures class, from 11:00 to 12:50 you have Elements of computing class, from 1:00 to 2:00 you have brake, and today you have Electronics lab from 2:00",
        "saturday" : "Boss today you have a full day of classes from 9:00 to 11:50 you have maths lab, from 12:00 to 12:50 you have english class, from 1:00 to 2:00 you have brake, and today you have elements of computing lab from 2:00",
        "sunday":"Boss today is holiday but we can't say anything when they will bomb with any assisgnments"}
        if day in Week.keys():
            self.talk(Week[day])

    #college resources commands
    #links
    def college(self,command):
        print(command)
        if 'open teams' in command:
            self.speak('opening your microsoft teams')
            webbrowser.open('https://teams.microsoft.com/')
        elif 'open lms' in command:
            self.speak('opening your lms sir')
            webbrowser.open('https://lms.mlrit.ac.in/?redirect=0')
        elif 'open mlrit portal' in command:
            self.speak('opening your mlrit portal')
            webbrowser.open('https://mlrit.ac.in/')
        else :
            self.No_result_found()
    
    #Online classes
    def OnlineClasses(self,command):
        print(command)
        #Keep as many "elif" statemets based on your subject Eg: I have kept a dummy links for JAVA and mechanics classes link of MS Teams
        if "java" in command:
            self.speak('opening DSA class in teams')
            webbrowser.open("https://teams.microsoft.com/java")
        elif "mechanics" in command:
            self.speak('opening mechanics class in teams')
            webbrowser.open("https://teams.microsoft.com/mechanics")
        elif 'online classes' in command:
            self.speak('opening your microsoft teams')
            webbrowser.open('https://teams.microsoft.com/')

    #Brower Search commands
    def B_S(self,command):
        print(command)
        try:
            # ('what is meant by' in self.command) or ('tell me about' in self.command) or ('who the heck is' in self.command)
            if 'wikipedia' in command:
                target1 = command.replace('search for','')
                target1 = target1.replace('in wikipedia','')
            elif 'what is meant by' in command:
                target1 = command.replace("what is meant by"," ")
            elif 'tell me about' in command:
                target1 = command.replace("tell me about"," ")
            elif 'who the heck is' in command :
                target1 = command.replace("who the heck is"," ")
            print("searching....")
            info = wikipedia.summary(target1,5)
            print(info)
            self.speak("according to wikipedia "+info)
        except :
            self.No_result_found()
        
    #Browser
    def brows(self,command):
        print(command)
        if 'google' in command:
            self.speak("Boss, what should I search on google..")
            S = self.take_command_with_fallback()#taking command for what to search in google
            webbrowser.open(f"{S}")
        elif 'edge' in command:
            self.speak('opening your Miscrosoft edge')
            os.startfile('..\\..\\MicrosoftEdge.exe')#path for your edge browser application
        else :
            self.No_result_found()

    #google applications selection
    #if there is any wrong with the URL's replace them with your browsers URL's
    def Google_Apps(self,command):
        print(command)
        if 'gmail' in command:
            self.speak('opening your google gmail')
            webbrowser.open('https://mail.google.com/mail/')
        elif 'maps' in command:
            self.speak('opening google maps')
            webbrowser.open('https://www.google.co.in/maps/')
        elif 'news' in command:
            self.speak('opening google news')
            webbrowser.open('https://news.google.com/')
        elif 'calender' in command:
            self.speak('opening google calender')
            webbrowser.open('https://calendar.google.com/calendar/')
        elif 'photos' in command:
            self.speak('opening your google photos')
            webbrowser.open('https://photos.google.com/')
        elif 'documents' in command:
            self.speak('opening your google documents')
            webbrowser.open('https://docs.google.com/document/')
        elif 'spreadsheet' in command:
            self.speak('opening your google spreadsheet')
            webbrowser.open('https://docs.google.com/spreadsheets/')
        else :
            self.No_result_found()


    def windows_manager(self,command):
        print(command)
        if "maximize this window" in command:
            pyautogui.hotkey('alt', 'space')
            time.sleep(1)
            pyautogui.press('x')

        elif "minimize the window" in command:
            pyautogui.hotkey('alt', 'space')
            time.sleep(1)
            pyautogui.press('n')

        elif 'switch window' in command:
            pyautogui.hotkey('alt', 'tab')

        elif "open new window" in command:
            pyautogui.hotkey('ctrl', 'n')

        elif 'open new tab' in command:
            pyautogui.hotkey('ctrl', 't')

        elif 'close current tab' in command:
            pyautogui.hotkey('ctrl', 'w')

        elif 'keep windows side by side' in command:
            self.arrange_windows_side_by_side()

        elif 'go to next window' in command:
            pyautogui.hotkey('alt', 'tab', interval=0.5)

        elif 'go to previous window' in command:
            pyautogui.hotkey('alt', 'shift', 'tab', interval=0.5)
        
        elif 'open history' in command:
            pyautogui.hotkey('ctrl', 'h')

        elif 'open downloads' in command:
            pyautogui.hotkey('ctrl', 'j')

        elif 'clear browsing history' in command:
            pyautogui.hotkey('ctrl', 'shift', 'delete')

        elif 'close all tabs' in command:
            pyautogui.hotkey('ctrl', 'shift', 'w')
        
            
    #youtube
    def yt(self,command):
        print(command)
        if 'play' in command:
            self.speak(" can you please say the name of the song")
            song = self.take_Command()
            if "play" in song:
                song = song.replace("play","")
            self.speak('playing '+song)
            print(f'playing {song}')
            pywhatkit.playonyt(song)
            print('playing')
        elif "download" in command:
            self.speak(" please enter the youtube video link which you want to download")
            link =self.take_command_with_fallback()
            yt=YouTube(link)
            yt.streams.get_highest_resolution().download()
            self.speak(f"Boss downloaded {yt.title} from the link you given into the main folder")
        elif "open youtube" in command:
            self.speak("Opening YouTube sir...")
            self.speak("What would you like to watch sir?")
            qrry = self.take_command_with_fallback().lower()
            webbrowser.open(f"https://www.youtube.com/results?search_query={qrry}")
        else :
            self.No_result_found()
        
    #Opensource accounts
    def open_source(self,command):
        print(command)
        if 'github' in command:
            self.speak('opening your github')
            webbrowser.open('https://github.com/BolisettySujith')
        elif 'gitlab' in command:
            self.speak('opening your gitlab')
            webbrowser.open('https://gitlab.com/-/profile')
        else :
            self.No_result_found()

    #Photo shops
    def edit(self,command):
        print(command)
        if 'slides' in command:
            self.speak('opening your google slides')
            webbrowser.open('https://docs.google.com/presentation/')
        elif 'canva' in command:
            self.speak('opening your canva')
            webbrowser.open('https://www.canva.com/')
        else :
            self.No_result_found()

    #OTT 
    def OTT(self,command):
        print(command)
        if 'hotstar' in command:
            self.speak('opening your disney plus hotstar')
            webbrowser.open('https://www.hotstar.com/in')
        elif 'prime' in command:
            self.speak('opening your amazon prime videos')
            webbrowser.open('https://www.primevideo.com/')
        elif 'netflix' in command:
            self.speak('opening Netflix videos')
            webbrowser.open('https://www.netflix.com/')
        else :
            self.No_result_found()

    #PC allications
    #NOTE: place the correct path for the applications from your PC there may be some path errors so please check the applications places
    #if you don't have any mentioned applications delete the codes for that
    #I have placed applications path based on my PC path check while using which OS you are using and change according to it
    def OpenApp(self,command):
        print(command)
        if ('calculator'in command) :
            self.speak('Opening calculator')
            os.startfile('C:\\Windows\\System32\\calc.exe')
        elif ('paint'in command) :
            self.speak('Opening msPaint')
            os.startfile('c:\\Windows\\System32\\mspaint.exe')
        elif ('notepad'in command) :
            self.speak('Opening notepad')
            os.startfile('c:\\Windows\\System32\\notepad.exe')
        elif ('discord'in command) :
            self.speak('Opening discord')
            os.startfile('..\\..\\Discord.exe')
        elif ('editor'in command) :
            self.speak('Opening your Visual studio code')
            os.startfile('..\\..\\Code.exe')
        elif ('online classes'in command) :
            self.speak('Opening your Microsoft teams')
            webbrowser.open('https://teams.microsoft.com/')
        elif ('spotify'in command) :
            self.speak('Opening spotify')
            os.startfile('..\\..\\Spotify.exe')
        elif ('media player'in command) :
            self.speak('Opening VLC media player')
            os.startfile("C:\Program Files\VideoLAN\VLC\vlc.exe")
        else :
            self.No_result_found()
            
    #closeapplications function
    def CloseApp(self,command):
        print(command)
        if ('calculator'in command) :
            self.speak("okay boss, closeing caliculator")
            os.system("taskkill /f /im calc.exe")
        elif ('paint'in command) :
            self.speak("okay boss, closeing mspaint")
            os.system("taskkill /f /im mspaint.exe")
        elif ('notepad'in command) :
            self.speak("okay boss, closeing notepad")
            os.system("taskkill /f /im notepad.exe")
        elif ('discord'in command) :
            self.speak("okay boss, closeing discord")
            os.system("taskkill /f /im Discord.exe")
        elif ('editor'in command) :
            self.speak("okay boss, closeing vs code")
            os.system("taskkill /f /im Code.exe")
        elif ('spotify'in command) :
            self.speak("okay boss, closeing spotify")
            os.system("taskkill /f /im Spotify.exe")
        elif ('media player'in command) :
            self.speak("okay boss, closeing media player")
            os.system("taskkill /f /im vlc.exe")
        else :
            self.No_result_found()

    #Shopping links
    def shopping(self,command):
        print(command)
        if 'flipkart' in command:
            self.speak('Opening flipkart online shopping website')
            webbrowser.open("https://www.flipkart.com/")
        elif 'amazon' in command:
            self.speak('Opening amazon online shopping website')
            webbrowser.open("https://www.amazon.in/")
        else :
            self.No_result_found()

    #PDF reader
    def search_file(self,file_name, directory, extensions):
        for root, dirs, files in os.walk(directory):
            for file in files:
                if file_name.lower() in file.lower() and any(file.lower().endswith(ext) for ext in extensions):
                    return os.path.join(root, file)
        return None

    def open_pdf(self,search_file):
        self.speak("Please tell me the name of the PDF file.")
        file_name = self.take_command_with_fallback()
        if file_name == "none":
            return
        home_dir = os.path.expanduser("~")
        file_path = search_file(file_name, home_dir, [".pdf"])
        if file_path:
            self.speak(f"Opening {file_path}")
            os.startfile(file_path)
        else:
            self.speak("File not found. Please try again.")

    def pdf_to_word(self,output_folder,search_file):
        self.speak("Please tell me the name of the PDF file to convert to Word.")
        file_name = self.take_command_with_fallback()
        if file_name == "none":
            return
        home_dir = os.path.expanduser("~")
        pdf_path = search_file(file_name, home_dir, [".pdf"])
        if not pdf_path:
            self.speak("PDF file not found. Please try again.")
            return
        try:
            docx_path = os.path.join(output_folder, "converted_document.docx")
            pdf2docx.parse(pdf_path, docx_path)
            self.speak("PDF has been converted to Word document successfully.")
        except Exception as ex:
            self.speak(f"An error occurred while converting PDF to Word document: {ex}")

    def pdf_to_ppt(self,output_folder,search_file):
        self.speak("Please tell me the name of the PDF file to convert to PowerPoint.")
        file_name = self.take_command_with_fallback()
        if file_name == "none":
            return
        home_dir = os.path.expanduser("~")
        pdf_path = search_file(file_name, home_dir, [".pdf"])
        if not pdf_path:
            self.speak("PDF file not found. Please try again.")
            return
        try:
            ppt_path = os.path.join(output_folder, "converted_presentation.pptx")
            pdf_reader = PyPDF2.PdfFileReader(pdf_path)
            ppt = Presentation()
            for page_num in range(pdf_reader.getNumPages()):
                page = pdf_reader.getPage(page_num)
                text = page.extract_text()
                slide_layout = ppt.slide_layouts[1]
                slide = ppt.slides.add_slide(slide_layout)
                textbox = slide.shapes.title
                textbox.text = text
            ppt.save(ppt_path)
            self.speak("PDF has been converted to PowerPoint presentation successfully.")
        except Exception as ex:
            self.speak(f"An error occurred while converting PDF to PowerPoint presentation: {ex}")

    def word_to_pdf(self,search_file,output_folder):
        self.speak("Please tell me the name of the Word document to convert to PDF.")
        file_name = self.take_command_with_fallback()
        if file_name == "none":
            return
        home_dir = os.path.expanduser("~")
        docx_path = search_file(file_name, home_dir, [".docx"])
        if not docx_path:
            self.speak("Word document not found. Please try again.")
            return
        try:
            pdf_path = os.path.join(output_folder, "converted_document.pdf")
            docx2pdf_convert(docx_path, pdf_path)
            self.speak("Word document has been converted to PDF successfully.")
        except Exception as ex:
            self.speak(f"An error occurred while converting Word document to PDF: {ex}")

    def ppt_to_pdf(self,search_file,output_folder):
        self.speak("Please tell me the name of the PowerPoint presentation to convert to PDF.")
        file_name = self.take_command_with_fallback()
        if file_name == "none":
            return
        home_dir = os.path.expanduser("~")
        ppt_path = search_file(file_name, home_dir, [".pptx"])
        if not ppt_path:
            self.speak("PowerPoint presentation not found. Please try again.")
            return
        try:
            pdf_path = os.path.join(output_folder, "converted_presentation.pdf")
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            ppt = powerpoint.Presentations.Open(ppt_path)
            ppt.SaveAs(pdf_path, 32)  # 32 for ppt to pdf
            ppt.Close()
            powerpoint.Quit()
            self.speak("PowerPoint presentation has been converted to PDF successfully.")
        except Exception as ex:
            self.speak(f"An error occurred while converting PowerPoint presentation to PDF: {ex}")

    def documents_manager(self,command,output_folder,search_file):
        print(command)
        if "open pdf" in command:
            self.open_pdf(search_file)
        elif "convert pdf to word" in command:
            self.pdf_to_word(output_folder,search_file)
        elif "convert pdf to ppt" in command:
            self.pdf_to_ppt(output_folder,search_file)
        elif "convert word to pdf" in command:
            self.word_to_pdf(output_folder,search_file)
        elif "convert ppt to pdf" in command:
            self.ppt_to_pdf(output_folder,search_file)


    #Time caliculating algorithm
    def silenceTime(self,command):
        print(command)
        x=0
        #caliculating the given time to seconds from the speech commnd string
        if ('10' in command) or ('ten' in command):x=600
        elif '1' in command or ('one' in command):x=60
        elif '2' in command or ('two' in command):x=120
        elif '3' in command or ('three' in command):x=180
        elif '4' in command or ('four' in command):x=240
        elif '5' in command or ('five' in command):x=300
        elif '6' in command or ('six' in command):x=360
        elif '7' in command or ('seven' in command):x=420
        elif '8' in command or ('eight' in command):x=480
        elif '9' in command or ('nine' in command):x=540
        self.silence(x)
        
    #Silence
    def silence(self,k):
        t = k
        s = "I will be silent for "+str(t/60)+" minutes"
        self.speak(s)
        while t:
            mins, secs = divmod(t, 60)
            timer = '{:02d}:{:02d}'.format(mins, secs)
            print(timer, end="\r")
            time.sleep(1)
            t -= 1
        self.talk("Boss "+str(k/60)+" minutes over")


    #location
    def locaiton(self):
        self.speak("Wait boss, let me check")
        try:
            IP_Address = get('https://api.ipify.org').text
            print(IP_Address)
            url = 'https://get.geojs.io/v1/ip/geo/'+IP_Address+'.json'
            print(url)
            geo_reqeust = get(url)
            geo_data = geo_reqeust.json()
            city = geo_data['city']
            state = geo_data['region']
            country = geo_data['country']
            tZ = geo_data['timezone']
            longitude = geo_data['longitude']
            latidute = geo_data['latitude']
            org = geo_data['organization_name']
            print(city+" "+state+" "+country+" "+tZ+" "+longitude+" "+latidute+" "+org)
            self.speak(f"Boss i am not sure, but i think we are in {city} city of {state} state of {country} country")
            self.speak(f"and boss, we are in {tZ} timezone the latitude os our location is {latidute}, and the longitude of our location is {longitude}, and we are using {org}\'s network ")
        except Exception as e:
            self.speak("Sorry boss, due to network issue i am not able to find where we are.")
            pass

    #ScreenShot
    def scshot(self):
        self.speak(" please tell me the name for this screenshot file")
        name = self.take_command_with_fallback()
        self.speak("Please hold the screen for few seconds, I am taking screenshot")
        time.sleep(3)
        img = pyautogui.screenshot()
        img.save(f"{name}.png")
        self.speak("he screenshot is saved in main folder.")

    #Arranging windows side bu side
    def arrange_windows_side_by_side(self):
        screen_width, screen_height = pyautogui.size()
    
    # Example positions and sizes for two windows
        window1_x = 0
        window1_y = 0
        window1_width = screen_width // 2
        window1_height = screen_height
    
        window2_x = screen_width // 2
        window2_y = 0
        window2_width = screen_width // 2
        window2_height = screen_height
    
    # Move and resize windows
        pyautogui.moveTo(window1_x + window1_width // 2, window1_y + window1_height // 2)
        pyautogui.hotkey('winleft', 'left')
        time.sleep(0.5)

        pyautogui.moveTo(window2_x + window2_width // 2, window2_y + window2_height // 2)
        pyautogui.hotkey('winleft', 'right')
        time.sleep(0.5)


    #News
    def news(self):
        self.speak("Today's news is")
        googlenews = GoogleNews(period='1d')
        googlenews.search("India")
        result = googlenews.result()
        data = pd.DataFrame(result)
        for i in range(len(data)):
            headline = data.iloc[i]['title']
            print(headline)
            self.speak(headline)

    #System condition
    def condition(self):
        usage = str(psutil.cpu_percent())
        self.speak("CPU is at"+usage+" percentage")
        battray = psutil.sensors_battery()
        percentage = battray.percent
        self.speak(f" our system have {percentage} percentage Battery")
        if percentage >=75:
            self.speak(f" we could have enough charging to continue our work")
        elif percentage >=40 and percentage <=75:
            self.speak(f" we should connect out system to charging point to charge our battery")
        elif percentage >=15 and percentage <=30:
            self.speak(f" we don't have enough power to work, please connect to charging")
        else:
            self.speak(f" we have very low power, please connect to charging otherwise the system will shutdown very soon")

    #no result found          
    def No_result_found(self):
        self.speak('Boss I couldn\'t understand, could you please say it again.')

if __name__ == "__main__":
    output_folder = "C:/Users/user/Desktop"
    jarvis = Jarvis()
    jarvis.run_jarvis("output_folder")
